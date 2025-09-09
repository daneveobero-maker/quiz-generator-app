import streamlit as st
import os
import docx
import pptx
import json
import re
from openpyxl import Workbook, load_workbook
import math
import PyPDF2
import io

# Placeholder for Gemini API
try:
    import google.generativeai as genai
except ImportError:
    st.error("The 'google-generativeai' library is not installed. Please install it using 'pip install google-generativeai'")
    genai = None

# --- CORE FUNCTIONS ---
def clean_json_response(response_text):
    if isinstance(response_text, (list, dict)): return response_text
    if isinstance(response_text, str):
        cleaned = re.sub(r"^```[a-zA-Z]*\n", "", response_text.strip())
        cleaned = re.sub(r"```$", "", cleaned.strip())
        try:
            return json.loads(cleaned)
        except json.JSONDecodeError as e:
            st.warning(f"JSON parsing failed. Error: {e}")
            return None
    st.error(f"Unexpected response type from API: {type(response_text)}")
    return None

def extract_text(uploaded_file):
    text_content = ""
    file_name = uploaded_file.name
    try:
        if file_name.endswith(".docx"):
            doc = docx.Document(uploaded_file)
            text_content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif file_name.endswith(".pptx"):
            prs = pptx.Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif file_name.endswith(".pdf"):
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
        elif file_name.endswith(".xlsx"):
            wb = load_workbook(uploaded_file, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            text_content += str(cell.value) + " "
            text_content = text_content.strip() + "\n"
    except Exception as e:
        st.error(f"Failed to extract text from {file_name}: {e}")
        return ""
    return text_content.strip()

def get_exam_period(filename, term, summer_last_week_str):
    match = re.search(r'week (\d+)', filename, re.IGNORECASE)
    if not match: return ""
    week_num = int(match.group(1))
    if term == "Semester":
        if 1 <= week_num <= 5: return "Prelim"
        elif 7 <= week_num <= 11: return "Midterm"
        elif 13 <= week_num <= 17: return "Finals"
        else: return ""
    elif term == "Summer":
        try:
            last_midterm_week = int(re.search(r'(\d+)', summer_last_week_str).group(1))
            if 1 <= week_num <= last_midterm_week: return "Midterm"
            elif (last_midterm_week + 1) <= week_num <= 17: return "Finals"
            else: return ""
        except (AttributeError, ValueError): return ""
    return ""

def generate_prompt(q_type, mc_subtype, situational_text, difficulty, num_items, words_per_choice, body_text):
    # This function's logic remains the same
    if q_type == "Multiple Choice":
        subtype_rules = {
            "Best Answer": "Create standard best-answer multiple choice questions.",
            "Negative": "At least one option should be worded as an incorrect/NOT statement.",
            "Complex": "Questions should be scenario-based or case-type requiring higher-order thinking.",
            "Analogy": "Questions should be in analogy format (A is to B as C is to ?).",
            "Assertion-Reason": "Use Assertion (A) and Reason (R) type format, where the answer explains their relationship.",
            "Problem Solving": "Questions should involve problem-solving or calculations based on the provided text."
        }
        extra_rule = subtype_rules.get(mc_subtype, "Create standard best-answer multiple choice questions.")
        prompt = f"""
        Create {num_items} {situational_text} {q_type} quiz questions ({mc_subtype} type) from the following text.
        Difficulty: {difficulty}. {extra_rule}
        Return ONLY a valid JSON array of objects, with no explanations.
        Each object must have these keys: "Subject Code", "Description", "Exam Period", "Learning Outcome", "Topic", "QUESTION", "Type of Question", "Taxonomy", "Choice A", "Choice B", "Choice C", "Choice D", "Answer".
        ⚠️ Rules: Each choice must be approx {words_per_choice} words. Answer must be only the letter. No "all/none of the above". "Type of Question" must be "Multiple Choice". "Taxonomy" must be one of: [Remembering, Understanding, Applying, Analyzing, Evaluating, Creating]. Subject code, Exam Period, Description must be empty strings.
        Text source: {body_text}
        """
    else:
        prompt = f"""
        Create {num_items} {situational_text} {q_type} quiz questions from the following text.
        Difficulty: {difficulty}.
        Return ONLY a valid JSON array of objects, with no explanations.
        Each object must have these keys: "Subject Code", "Description", "Exam Period", "Learning Outcome", "Topic", "QUESTION", "Type of Question", "Taxonomy", "Answer".
        ⚠️ Rules: Do NOT include multiple-choice options. For "True or False", Answer must be "True" or "False". For "Short Answer", integrate a quote instruction.
        Text source: {body_text}
        """
    return prompt

def call_gemini(api_key, model, prompt):
    if not genai:
         st.error("Gemini API library is not available.")
         return None
    try:
        genai.configure(api_key=api_key)
        llm = genai.GenerativeModel(model)
        response = llm.generate_content(prompt)
        return response.text
    except Exception as e:
        st.error(f"An error occurred with the Gemini API call: {e}")
        return None

def save_to_word(quiz_data):
    doc = docx.Document()
    doc.add_heading("Generated Quiz", level=1)
    for i, q in enumerate(quiz_data, 1):
        doc.add_paragraph(f"{i}. {q.get('QUESTION', '')}", style="List Number")
        for opt in ["Choice A", "Choice B", "Choice C", "Choice D"]:
            if q.get(opt): doc.add_paragraph(f"    {opt[-1]}. {q[opt]}")
        p = doc.add_paragraph()
        p.add_run(f"Answer: {q.get('Answer', '')}").bold = True
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def save_to_excel(quiz_data, q_type):
    wb = Workbook()
    ws = wb.active
    ws.title = "Quiz Questions"
    if q_type == "Multiple Choice":
        headers = ["Subject Code", "Description", "Exam Period", "Learning Outcome", "Topic", "QUESTION", "Type of Question", "Taxonomy", "Choice A", "Choice B", "Choice C", "Choice D", "Answer"]
    else:
        headers = ["Subject Code", "Description", "Exam Period", "Learning Outcome", "Topic", "QUESTION", "Type of Question", "Taxonomy", "Answer"]
    ws.append(headers)
    for q in quiz_data:
        row = [q.get(h, "") for h in headers]
        ws.append(row)
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer

# --- Streamlit UI ---
st.set_page_config(page_title="Quiz Generator with Gemini", layout="wide")
st.title("📚 Quiz Generator with Gemini")

# --- Sidebar for Configuration ---
with st.sidebar:
    st.header("⚙️ Configuration")
    # Try to get the API key from Streamlit's secrets, fallback to text input
    try:
        api_key = st.secrets["GEMINI_API_KEY"]
    except (KeyError, FileNotFoundError):
        api_key = st.text_input("Enter your Gemini API Key", type="password")
    
    model = st.selectbox("Select Gemini Model", ["gemini-2.5-flash", "gemini-2.0-flash"])
    q_type = st.selectbox("Select Quiz Type", ["Multiple Choice", "Identification", "Short Answer", "True or False", "Modified True or False"])
    difficulty = st.selectbox("Select Difficulty", ["Easy", "Medium", "Hard"])
    term = st.selectbox("Select Term", ["Semester", "Summer"])
    summer_last_week = None
    if term == "Summer":
        summer_last_week = st.selectbox("Midterm Last Week (Summer)", ["Week 7", "Week 8", "Week 9", "Week 10"])
    total_items_per_file = st.number_input("Number of Items per File (Max 50)", 1, 50, 10)
    words_per_choice = st.number_input("Number of words per choice (for MC)", 5, 50, 10)
    is_situational = st.checkbox("Generate Situational Questions")

    mc_subtypes_selected = []
    if q_type == "Multiple Choice":
        with st.expander("Multiple Choice Subtypes"):
            num_loops = math.ceil(total_items_per_file / 10)
            mc_options = ["Best Answer", "Negative", "Complex", "Analogy", "Assertion-Reason", "Problem Solving"]
            for i in range(num_loops):
                subtype = st.selectbox(f"Batch {i+1} Subtype", mc_options, index=(i % len(mc_options)), key=f"mc_subtype_{i}")
                mc_subtypes_selected.append(subtype)

# --- Main Page ---
st.info("Upload one or more Word, PowerPoint, PDF, or Excel files to generate a quiz.")
uploaded_files = st.file_uploader(
    "Select Files",
    type=['docx', 'pptx', 'pdf', 'xlsx'],
    accept_multiple_files=True
)

if st.button("🚀 Generate Quiz", disabled=(not uploaded_files)):
    if not api_key: 
        st.error("Please enter your Gemini API Key in the sidebar.")
        st.stop()
    
    all_quiz_data = []
    num_loops_per_file = math.ceil(total_items_per_file / 10)
    total_api_calls = num_loops_per_file * len(uploaded_files)
    progress_bar = st.progress(0)
    current_api_call = 0
    with st.status("Generating Quiz...", expanded=True) as status:
        for i, file in enumerate(uploaded_files):
            st.write(f"**Processing File {i+1}/{len(uploaded_files)}: `{file.name}`**")
            body_text = extract_text(file)
            if not body_text:
                st.warning(f"Skipping `{file.name}` as no text could be extracted.")
                current_api_call += num_loops_per_file
                progress_bar.progress(current_api_call / total_api_calls if total_api_calls > 0 else 0)
                continue

            topic_name = os.path.splitext(file.name)[0]
            exam_period = get_exam_period(file.name, term, summer_last_week)

            for loop_num in range(num_loops_per_file):
                remaining_questions = total_items_per_file - (loop_num * 10)
                questions_this_loop = min(10, remaining_questions)
                st.write(f"  - Calling API for batch {loop_num + 1}/{num_loops_per_file}...")
                
                mc_subtype = mc_subtypes_selected[loop_num] if q_type == "Multiple Choice" else None
                situational_text = "situational" if is_situational else ""
                
                prompt = generate_prompt(q_type, mc_subtype, situational_text, difficulty, questions_this_loop, words_per_choice, body_text)
                response_text = call_gemini(api_key, model, prompt)
                quiz_data_batch = clean_json_response(response_text)
                
                current_api_call += 1
                progress_bar.progress(current_api_call / total_api_calls if total_api_calls > 0 else 0)

                if quiz_data_batch and isinstance(quiz_data_batch, list):
                    for question in quiz_data_batch:
                        question['Topic'] = topic_name
                        question['Exam Period'] = exam_period
                    all_quiz_data.extend(quiz_data_batch)
                else:
                    st.warning(f"Batch {loop_num + 1} for `{file.name}` returned invalid data. Skipping batch.")
        
        status.update(label="Quiz generation complete!", state="complete")

    # --- Download Buttons ---
    if all_quiz_data:
        st.success(f"Successfully generated a total of **{len(all_quiz_data)}** questions!")
        
        excel_buffer = save_to_excel(all_quiz_data, q_type)
        word_buffer = save_to_word(all_quiz_data)
        
        col1, col2 = st.columns(2)
        with col1:
            st.download_button(
                label="📥 Download as Excel (.xlsx)",
                data=excel_buffer,
                file_name="Generated_Quiz.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
        with col2:
            st.download_button(
                label="📥 Download as Word (.docx)",
                data=word_buffer,
                file_name="Generated_Quiz.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
    else:
        st.error("No valid quiz data was generated from any of the files.")

